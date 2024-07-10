"""Module managing base 64 documents"""

import csv
import base64
from io import BytesIO
from io import StringIO
from typing import List, Optional, Dict

from pydantic import BaseModel

from langchain.schema import Document

import fitz
import openpyxl
from docx import Document as DocxDocument

from docx.text.paragraph import Paragraph
from docx.table import Table

from pptx import Presentation

from modules import logger

from config.config import EnvParam


class Documents(BaseModel):
    """Class with all documents"""

    base64_documents: Optional[List[Dict[str, str]]]

    documents_list: list[Document] = []

    documents_names: str = ""

    def create_documents(self) -> None:
        """Create all documents to langchain docs"""

        if self.base64_documents:
            for document in self.base64_documents:
                new_document: Optional[Document] = self.create_document(
                    documents=document
                )

                if new_document:
                    self.documents_list.append(new_document)

                    self.documents_names += f"- {new_document.metadata['file_name']}\n"

        logger.info(f"User documents added:\n{self.documents_names}")

    def create_document(self, documents: Dict[str, str]) -> Optional[Document]:
        """Create one document"""

        try:
            base64_content: Optional[str] = documents.get(EnvParam.BASE64_PARAM_NAME)

            metadata: Optional[str] = documents.get(EnvParam.DOCUMENT_METADATA_NAME)

            if not base64_content:
                return None

            if not metadata:
                metadata = "Any file"

            content: Optional[str] = self.base_64_to_txt(
                base64_content=base64_content, metadata=metadata
            )

            if content:
                return Document(page_content=content, metadata={"file_name": metadata})

            logger.warning("Document %s is empty", metadata)

        except Exception as e:
            logger.error("Error creating document : %s", e)

        return None

    def base_64_to_txt(self, base64_content: str, metadata: str) -> Optional[str]:
        """Convert base 64 to text from .txt, .pdf, and .docx files."""

        try:
            _, encoded_data = base64_content.split(",", 1)

            decoded_data = base64.b64decode(encoded_data)

            extension = "." + metadata.split(".")[-1]

            logger.info("File found %s (%s)", metadata, extension)

            text_content: Optional[str] = self.extract_content(
                decoded_data=decoded_data, extension=extension
            )

            return text_content

        except Exception as e:
            logger.error("Error getting text from base64 %s", e)

            return None

    def extract_content(
        self, decoded_data: bytes, extension: Optional[str]
    ) -> Optional[str]:
        """Extract content from bytes according file extension"""

        if extension in [".txt"]:
            text_content: Optional[str] = decoded_data.decode("utf-8")

        elif extension in [".pdf"]:
            text_content = self.extract_text_from_pdf(decoded_data)

        elif extension in [".docx"]:
            text_content = self.extract_text_from_docx(decoded_data)

        elif extension == ".pptx":
            text_content = self.extract_text_from_pptx(decoded_data)

        elif extension == ".xlsx":
            text_content = self.extract_text_from_xlsx(decoded_data)

        elif extension == ".csv":
            text_content = self.extract_text_from_csv(decoded_data)

        else:
            logger.error("Unsupported file type: %s", extension)

            text_content = None

        return text_content

    def extract_text_from_pdf(self, pdf_bytes) -> Optional[str]:
        """Convert pdf to text"""

        try:
            text = ""

            with fitz.open(stream=pdf_bytes, filetype="pdf") as doc:
                for page in doc:
                    text += page.get_text()

        except Exception as e:
            logger.error("Error converting base64 to pdf %s", e)

            text = None

        logger.info("PDF convert succesfully")

        return text

    def extract_text_from_docx(self, docx_bytes) -> Optional[str]:
        """Convert docx to text"""

        def clean_val(val: str) -> str:
            return val.strip().replace("\n", "").replace("\t", "")

        try:
            doc = DocxDocument(BytesIO(docx_bytes))

            full_content = []

            for element in doc.iter_inner_content():
                if isinstance(element, Paragraph):
                    full_content.append(f"{element.text}")
                elif isinstance(element, Table):
                    table_text = "\n".join(
                        [
                            "; ".join(
                                f"{clean_val(element.rows[0].cells[i].text)}: {clean_val(cell.text)}"
                                for i, cell in enumerate(row.cells)
                            )
                            for row in element.rows[1:]
                        ]
                    )
                    full_content.append(table_text)

            text = "\n\n".join(full_content)

        except Exception as e:
            logger.error("Error converting base64 to docx %s", e)

            text = None

        logger.info("DOCX convert succesfully => %s", text)

        return text

    def extract_text_from_pptx(self, pptx_bytes) -> Optional[str]:
        """Convert pptx to text"""
        try:
            prs = Presentation(BytesIO(pptx_bytes))
            text = ""
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"):
                        text += shape.text + "\n"
        except Exception as e:
            logger.error("Error converting PPTX to text: %s", e)
            text = None
        else:
            logger.info("PPTX converted successfully.")
        return text

    def extract_text_from_xlsx(self, xlsx_bytes) -> Optional[str]:
        """Convert xlsx to text"""
        try:
            wb = openpyxl.load_workbook(BytesIO(xlsx_bytes), data_only=True)
            text = ""
            for sheet in wb:
                for row in sheet.iter_rows(values_only=True):
                    if row:
                        text += (
                            ",".join(
                                [str(cell) if cell is not None else "" for cell in row]
                            )
                            + "\n"
                        )
        except Exception as e:
            logger.error("Error converting XLSX to text: %s", e)
            text = None
        else:
            logger.info("XLSX converted successfully.")
        return text

    def extract_text_from_csv(
        self, csv_bytes: bytes, delimiter=";", row_delimiter="\n", encoding="utf-8"
    ) -> Optional[str]:
        """Convert CSV bytes to a sorted text representation"""
        text = None
        try:
            csv_file = StringIO(csv_bytes.decode(encoding))

            csv_reader = csv.reader(csv_file, delimiter=delimiter)

            header = next(csv_reader)

            data = list(csv_reader)

            data_sorted = sorted(data, key=lambda row: tuple(row))

            csv_list = data_sorted

            formatted_lines = []

            for row in csv_list:
                formatted_row = "; ".join(
                    f"{name}: {value}" for name, value in zip(header, row)
                )
                formatted_lines.append(formatted_row + ";")

            text = row_delimiter.join(formatted_lines)

            logger.error("CSV converted successfully => %s", text)

        except Exception as e:
            logger.error("Error converting CSV to text: %s", e)

        return text
